import os
import sqlite3
import datetime as dt
from pathlib import Path

import numpy as np
import pandas as pd
import yfinance as yf
from pptx import Presentation
from pptx.util import Inches, Pt

DB_PATH = Path(os.getenv("ETF_DB_PATH", "C:/Users/bobi/.openclaw/workspace/etf_analytics.db"))
PPT_OUTPUT_DIR = Path(os.getenv(
    "ETF_PPT_DIR",
    r"C:/Users/bobi/OneDrive/바탕 화면/OneDrive/gram to vivo/☆취업/★자산운용사 취업 포트폴리오",
))
PPT_NAME = os.getenv("ETF_PPT_NAME", "ETF 비교분석.pptx")

# 후보군에서 최근 30일 거래대금 상위 10개를 선정
ETF_UNIVERSE = [
    "SPY", "QQQ", "VTI", "IVV", "VOO", "IWM", "DIA", "EFA", "EEM", "XLF",
    "XLK", "XLV", "XLE", "XLI", "XLY", "XLP", "VNQ", "ARKK", "TLT", "GLD",
    "SLV", "HYG", "LQD", "SOXX", "SMH", "KWEB", "EWY", "VWO", "VEA", "SCHD",
]


def ensure_db(conn: sqlite3.Connection):
    conn.execute(
        """
        CREATE TABLE IF NOT EXISTS etf_daily (
            date TEXT NOT NULL,
            ticker TEXT NOT NULL,
            open REAL,
            high REAL,
            low REAL,
            close REAL,
            volume REAL,
            avg_dollar_vol_30d REAL,
            rsi14 REAL,
            ma20 REAL,
            ma60 REAL,
            macd REAL,
            macd_signal REAL,
            pe_ratio REAL,
            expense_ratio REAL,
            ytd_return REAL,
            PRIMARY KEY (date, ticker)
        )
        """
    )
    conn.commit()


def fetch_prices(tickers):
    df = yf.download(tickers, period="6mo", interval="1d", auto_adjust=False, progress=False)
    if isinstance(df.columns, pd.MultiIndex):
        return df
    raise RuntimeError("가격 데이터 형식이 예상과 다릅니다.")


def get_series(px_multi: pd.DataFrame, field: str, ticker: str) -> pd.Series:
    if (field, ticker) in px_multi.columns:
        s = px_multi[(field, ticker)].dropna()
        return s
    return pd.Series(dtype=float)


def calc_rsi(close: pd.Series, period=14):
    d = close.diff()
    up = d.clip(lower=0).rolling(period).mean()
    dn = (-d.clip(upper=0)).rolling(period).mean()
    rs = up / (dn + 1e-9)
    return 100 - 100 / (1 + rs)


def pick_top10_by_dollar_volume(px_multi: pd.DataFrame):
    rows = []
    for t in ETF_UNIVERSE:
        c = get_series(px_multi, "Close", t)
        v = get_series(px_multi, "Volume", t)
        if len(c) < 25 or len(v) < 25:
            continue
        dv = (c * v).tail(30).mean()
        rows.append((t, float(dv)))
    rows = sorted(rows, key=lambda x: x[1], reverse=True)
    return [x[0] for x in rows[:10]], {k: v for k, v in rows}


def collect_metrics(top10, px_multi, avg_dv_map):
    out = []
    asof = None

    for t in top10:
        o = get_series(px_multi, "Open", t)
        h = get_series(px_multi, "High", t)
        l = get_series(px_multi, "Low", t)
        c = get_series(px_multi, "Close", t)
        v = get_series(px_multi, "Volume", t)
        if len(c) < 70:
            continue

        asof = c.index[-1].strftime("%Y-%m-%d")
        rsi = calc_rsi(c).iloc[-1]
        ma20 = c.rolling(20).mean().iloc[-1]
        ma60 = c.rolling(60).mean().iloc[-1]
        ema12 = c.ewm(span=12, adjust=False).mean()
        ema26 = c.ewm(span=26, adjust=False).mean()
        macd = (ema12 - ema26).iloc[-1]
        macd_sig = (ema12 - ema26).ewm(span=9, adjust=False).mean().iloc[-1]

        info = yf.Ticker(t).info or {}
        pe = info.get("trailingPE")
        exp = info.get("annualReportExpenseRatio") or info.get("expenseRatio")
        ytd = info.get("ytdReturn")

        out.append(
            {
                "date": asof,
                "ticker": t,
                "open": float(o.iloc[-1]),
                "high": float(h.iloc[-1]),
                "low": float(l.iloc[-1]),
                "close": float(c.iloc[-1]),
                "volume": float(v.iloc[-1]),
                "avg_dollar_vol_30d": float(avg_dv_map.get(t, np.nan)),
                "rsi14": float(rsi),
                "ma20": float(ma20),
                "ma60": float(ma60),
                "macd": float(macd),
                "macd_signal": float(macd_sig),
                "pe_ratio": None if pe is None else float(pe),
                "expense_ratio": None if exp is None else float(exp),
                "ytd_return": None if ytd is None else float(ytd),
            }
        )

    return pd.DataFrame(out)


def upsert_db(conn: sqlite3.Connection, df: pd.DataFrame):
    if df.empty:
        return
    cols = [
        "date", "ticker", "open", "high", "low", "close", "volume", "avg_dollar_vol_30d",
        "rsi14", "ma20", "ma60", "macd", "macd_signal", "pe_ratio", "expense_ratio", "ytd_return",
    ]
    sql = """
    INSERT OR REPLACE INTO etf_daily
    (date, ticker, open, high, low, close, volume, avg_dollar_vol_30d,
     rsi14, ma20, ma60, macd, macd_signal, pe_ratio, expense_ratio, ytd_return)
    VALUES (?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?)
    """
    conn.executemany(sql, df[cols].itertuples(index=False, name=None))
    conn.commit()


def make_signal(row):
    s = 0
    if row["close"] > row["ma20"]:
        s += 1
    if row["ma20"] > row["ma60"]:
        s += 1
    if row["rsi14"] >= 65:
        s -= 1
    if row["rsi14"] <= 35:
        s += 1
    if row["macd"] > row["macd_signal"]:
        s += 1
    if s >= 2:
        return "상승 우위"
    if s <= 0:
        return "약세/중립"
    return "중립"


def build_ppt(df: pd.DataFrame):
    PPT_OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
    out_path = PPT_OUTPUT_DIR / PPT_NAME

    prs = Presentation()

    # Title slide
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "ETF 비교분석"
    subtitle = slide.placeholders[1]
    subtitle.text = f"자동 업데이트 리포트 | {dt.datetime.now().strftime('%Y-%m-%d %H:%M')}"

    # Summary slide
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "요약 (상위 10개 ETF)"
    tx = slide.shapes.add_textbox(Inches(0.6), Inches(1.3), Inches(12), Inches(5)).text_frame
    tx.word_wrap = True
    for _, r in df.head(10).iterrows():
        ytd = "N/A" if pd.isna(r["ytd_return"]) else f"{r['ytd_return']*100:.1f}%"
        pe = "N/A" if pd.isna(r["pe_ratio"]) else f"{r['pe_ratio']:.1f}"
        sig = make_signal(r)
        p = tx.add_paragraph()
        p.text = f"• {r['ticker']}: 종가 {r['close']:.2f}, RSI {r['rsi14']:.1f}, YTD {ytd}, PE {pe}, 판단 {sig}"
        p.font.size = Pt(16)

    # Table slide
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "기본/기술/가치 지표 비교"

    cols = ["ticker", "close", "rsi14", "ma20", "ma60", "pe_ratio", "expense_ratio", "ytd_return"]
    show = df[cols].copy().head(10)
    show["ytd_return"] = show["ytd_return"].apply(lambda x: np.nan if pd.isna(x) else x * 100)
    rows_n, cols_n = show.shape

    table = slide.shapes.add_table(rows_n + 1, cols_n, Inches(0.4), Inches(1.2), Inches(12.5), Inches(5.2)).table
    for j, c in enumerate(show.columns):
        table.cell(0, j).text = c
    for i in range(rows_n):
        for j, c in enumerate(show.columns):
            val = show.iloc[i, j]
            if isinstance(val, (float, np.floating)):
                if c == "ytd_return":
                    text = "N/A" if pd.isna(val) else f"{val:.1f}%"
                else:
                    text = "N/A" if pd.isna(val) else f"{val:.2f}"
            else:
                text = str(val)
            table.cell(i + 1, j).text = text

    prs.save(out_path)
    return out_path


def main():
    px = fetch_prices(ETF_UNIVERSE)
    top10, avg_dv_map = pick_top10_by_dollar_volume(px)
    df = collect_metrics(top10, px, avg_dv_map)
    if df.empty:
        raise RuntimeError("ETF 데이터 수집 결과가 비어 있습니다.")

    conn = sqlite3.connect(DB_PATH)
    try:
        ensure_db(conn)
        upsert_db(conn, df)
    finally:
        conn.close()

    ppt_path = build_ppt(df)
    print(f"DB updated: {DB_PATH}")
    print(f"PPT updated: {ppt_path}")


if __name__ == "__main__":
    main()
